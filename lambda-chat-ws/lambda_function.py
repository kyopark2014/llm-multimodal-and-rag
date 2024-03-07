import json
import boto3
import os
import time
import datetime
from io import BytesIO
import PyPDF2
import csv
import traceback
import re
from urllib import parse

from botocore.config import Config
from langchain.prompts import PromptTemplate
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains.summarize import load_summarize_chain
from langchain.chains import ConversationChain
from langchain.chains import LLMChain
from langchain.memory import ConversationBufferWindowMemory
from langchain.callbacks.streaming_stdout import StreamingStdOutCallbackHandler

from langchain_community.llms.bedrock import Bedrock
from langchain_community.docstore.document import Document
from langchain_community.vectorstores.faiss import FAISS
from langchain_community.vectorstores.opensearch_vector_search import OpenSearchVectorSearch
from langchain_community.embeddings import BedrockEmbeddings
from multiprocessing import Process, Pipe
from googleapiclient.discovery import build
from opensearchpy import OpenSearch
from langchain.schema import BaseMessage

from langchain_community.chat_models import BedrockChat
from langchain_core.messages import HumanMessage

s3 = boto3.client('s3')
s3_bucket = os.environ.get('s3_bucket') # bucket name
s3_prefix = os.environ.get('s3_prefix')
callLogTableName = os.environ.get('callLogTableName')

opensearch_account = os.environ.get('opensearch_account')
opensearch_passwd = os.environ.get('opensearch_passwd')
enableReference = os.environ.get('enableReference', 'false')
debugMessageMode = os.environ.get('debugMessageMode', 'false')
opensearch_url = os.environ.get('opensearch_url')
path = os.environ.get('path')
doc_prefix = s3_prefix+'/'
speech_prefix = 'speech/'

useParallelRAG = os.environ.get('useParallelRAG', 'true')
roleArn = os.environ.get('roleArn')
top_k = int(os.environ.get('numberOfRelevantDocs', '8'))
selected_LLM = 0
MSG_LENGTH = 100
MSG_HISTORY_LENGTH = 20
speech_generation = True
history_length = 0
token_counter_history = 0

minDocSimilarity = 200

# google search api
googleApiSecret = os.environ.get('googleApiSecret')
secretsmanager = boto3.client('secretsmanager')
try:
    get_secret_value_response = secretsmanager.get_secret_value(
        SecretId=googleApiSecret
    )
    #print('get_secret_value_response: ', get_secret_value_response)
    secret = json.loads(get_secret_value_response['SecretString'])
    #print('secret: ', secret)
    google_api_key = secret['google_api_key']
    google_cse_id = secret['google_cse_id']
    #print('google_cse_id: ', google_cse_id)    

except Exception as e:
    raise e

# websocket
connection_url = os.environ.get('connection_url')
client = boto3.client('apigatewaymanagementapi', endpoint_url=connection_url)
print('connection_url: ', connection_url)

HUMAN_PROMPT = "\n\nHuman:"
AI_PROMPT = "\n\nAssistant:"

map_chain = dict() # For RAG
map_chat = dict() # For general conversation  

# Multi-LLM
def get_llm(profile_of_LLMs, selected_LLM):
    profile = profile_of_LLMs[selected_LLM]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'LLM: {selected_LLM}, bedrock_region: {bedrock_region}, modelId: {modelId}')
    maxOutputTokens = int(profile['maxOutputTokens'])
                          
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }            
        )
    )
    parameters = {
        "max_tokens":maxOutputTokens,     
        "temperature":0.1,
        "top_k":250,
        "top_p":0.9,
        "stop_sequences": [HUMAN_PROMPT]
    }
    # print('parameters: ', parameters)

    llm = BedrockChat(
        model_id=modelId,
        client=boto3_bedrock, 
        streaming=True,
        callbacks=[StreamingStdOutCallbackHandler()],
        model_kwargs=parameters,
    )        
    return llm

def get_embedding(profile_of_LLMs, selected_LLM):
    profile = profile_of_LLMs[selected_LLM]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'LLM: {selected_LLM}, bedrock_region: {bedrock_region}, modelId: {modelId}')
    
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }            
        )
    )
    
    bedrock_embedding = BedrockEmbeddings(
        client=boto3_bedrock,
        region_name = bedrock_region,
        model_id = 'amazon.titan-embed-text-v1' 
    )  
    
    return bedrock_embedding

def sendMessage(id, body):
    try:
        client.post_to_connection(
            ConnectionId=id, 
            Data=json.dumps(body)
        )
    except Exception:
        err_msg = traceback.format_exc()
        print('err_msg: ', err_msg)
        raise Exception ("Not able to send a message")

def sendResultMessage(connectionId, requestId, msg):    
    result = {
        'request_id': requestId,
        'msg': msg,
        'status': 'completed'
    }
    #print('debug: ', json.dumps(debugMsg))
    sendMessage(connectionId, result)

def sendDebugMessage(connectionId, requestId, msg):
    debugMsg = {
        'request_id': requestId,
        'msg': msg,
        'status': 'debug'
    }
    #print('debug: ', json.dumps(debugMsg))
    sendMessage(connectionId, debugMsg)

def sendErrorMessage(connectionId, requestId, msg):
    errorMsg = {
        'request_id': requestId,
        'msg': msg,
        'status': 'error'
    }
    print('error: ', json.dumps(errorMsg))
    sendMessage(connectionId, errorMsg)

def isKorean(text):
    # check korean
    pattern_hangul = re.compile('[\u3131-\u3163\uac00-\ud7a3]+')
    word_kor = pattern_hangul.search(str(text))
    # print('word_kor: ', word_kor)

    if word_kor and word_kor != 'None':
        print('Korean: ', word_kor)
        return True
    else:
        print('Not Korean: ', word_kor)
        return False

def get_prompt_template(query, conv_type, rag_type):    
    if isKorean(query):
        if conv_type == "normal": # for General Conversation
            prompt_template = """\n\nHuman: 다음의 <history>는 Human과 Assistant의 친근한 이전 대화입니다. Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다.

            <history>
            {history}
            </history>            

            <question>            
            {input}
            </question>
            
            Assistant:"""

        elif conv_type=='qa':  
            prompt_template = """\n\nHuman: 다음의 <context> tag안의 참고자료를 이용하여 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다.
            
            <context>
            {context}
            </context>

            <question>            
            {question}
            </question>

            Assistant:"""
                
        elif conv_type == "translation":  # for translation, input
            prompt_template = """\n\nHuman: 다음의 <article>를 English로 번역하세요. 머리말은 건너뛰고 본론으로 바로 들어가주세요. 또한 결과는 <result> tag를 붙여주세요.

            <article>
            {input}
            </article>
                        
            Assistant:"""

        elif conv_type == "sentiment":  # for sentiment, input
            prompt_template = """\n\nHuman: 아래의 <example> review와 Extracted Topic and sentiment 인 <result>가 있습니다.
            <example>
            객실은 작지만 깨끗하고 편안합니다. 프론트 데스크는 정말 분주했고 체크인 줄도 길었지만, 직원들은 프로페셔널하고 매우 유쾌하게 각 사람을 응대했습니다. 우리는 다시 거기에 머물것입니다.
            </example>
            <result>
            청소: 긍정적, 
            서비스: 긍정적
            </result>

            아래의 <review>에 대해서 위의 <result> 예시처럼 Extracted Topic and sentiment 을 만들어 주세요..

            <review>
            {input}
            </review>

            Assistant:"""

        elif conv_type == "extraction":  # information extraction
            prompt_template = """\n\nHuman: 다음 텍스트에서 이메일 주소를 정확하게 복사하여 한 줄에 하나씩 적어주세요. 입력 텍스트에 정확하게 쓰여있는 이메일 주소만 적어주세요. 텍스트에 이메일 주소가 없다면, "N/A"라고 적어주세요. 또한 결과는 <result> tag를 붙여주세요.

            <text>
            {input}
            </text>

            Assistant:"""

        elif conv_type == "pii":  # removing PII(personally identifiable information) containing name, phone number, address
            prompt_template = """\n\nHuman: 아래의 <text>에서 개인식별정보(PII)를 모두 제거하여 외부 계약자와 안전하게 공유할 수 있도록 합니다. 이름, 전화번호, 주소, 이메일을 XXX로 대체합니다. 또한 결과는 <result> tag를 붙여주세요.
            
            <text>
            {input}
            </text>
        
            Assistant:"""

        elif conv_type == "grammar":  # Checking Grammatical Errors
            prompt_template = """\n\nHuman: 다음의 <article>에서 문장의 오류를 찾아서 설명하고, 오류가 수정된 문장을 답변 마지막에 추가하여 주세요.

            <article>
            {input}
            </article>
            
            Assistant: """

        elif conv_type == "step-by-step":  # compelex question 
            prompt_template = """\n\nHuman: 다음은 Human과 Assistant의 친근한 대화입니다. Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. 아래 문맥(context)을 참조했음에도 답을 알 수 없다면, 솔직히 모른다고 말합니다. 여기서 Assistant의 이름은 서연입니다.

            {input}

            Assistant: 단계별로 생각할까요?

            Human: 예, 그렇게하세요.
            
            Assistant:"""

        elif conv_type == "like-child":  # Child Conversation (few shot)
            prompt_template = """\n\nHuman: 다음 대화를 완성하기 위해 "A"로 말하는 다음 줄을 작성하세요. Assistant는 유치원 선생님처럼 대화를 합니다.
            
            Q: 이빨 요정은 실제로 있을까?
            A: 물론이죠, 오늘 밤 당신의 이를 감싸서 베개 밑에 넣어두세요. 아침에 뭔가 당신을 기다리고 있을지도 모릅니다.
            Q: {input}

            Assistant:"""      

        elif conv_type == "timestamp-extraction":
            prompt_template = """\n\nHuman: 아래의 <text>는 시간을 포함한 텍스트입니다. 친절한 AI Assistant로서 시간을 추출하여 아래를 참조하여 <example>과 같이 정리해주세요.
            
            - 년도를 추출해서 <year>/<year>로 넣을것 
            - 월을 추출해서 <month>/<month>로 넣을것
            - 일을 추출해서 <day>/<day>로 넣을것
            - 시간을 추출해서 24H으로 정리해서 <hour>/<hour>에 넣을것
            - 분을 추출해서 <minute>/<minute>로 넣을것

            이때의 예제는 아래와 같습니다.
            <example>
            2022년 11월 3일 18시 26분
            </example>
            <result>
                <year>2022</year>
                <month>11</month>
                <day>03</day>
                <hour>18</hour>
                <minute>26</minute>
            </result>

            결과에 개행문자인 "\n"과 글자 수와 같은 부가정보는 절대 포함하지 마세요.

            <text>
            {input}
            </text>

            Assistant:"""  

        elif conv_type == "funny": # for free conversation
            prompt_template = """\n\nHuman: 다음의 <history>는 Human과 Assistant의 친근한 이전 대화입니다. 모든 대화는 반말로하여야 합니다. Assistant의 이름은 서서이고 10살 여자 어린이 상상력이 풍부하고 재미있는 대화를 합니다. 때로는 바보같은 답변을 해서 재미있게 해줍니다.

            <history>
            {history}
            </history>

            <question>            
            {input}
            </question>
            
            Assistant:"""     

        elif conv_type == "get-weather":  # getting weather (function calling)
            prompt_template = """\n\nHuman: In this environment you have access to a set of tools you can use to answer the user's question.

            You may call them like this. Only invoke one function at a time and wait for the results before invoking another function:
            
            <function_calls>
            <invoke>
            <tool_name>$TOOL_NAME</tool_name>
            <parameters>
            <$PARAMETER_NAME>$PARAMETER_VALUE</$PARAMETER_NAME>
            ...
            </parameters>
            </invoke>
            </function_calls>

            Here are the tools available:
            <tools>
            {tools_string}
            </tools>

            Human:
            {user_input}

            Assistant:"""                  
                
        else:
            prompt_template = """\n\nHuman: 다음은 Human과 Assistant의 친근한 대화입니다. Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant는 모르는 질문을 받으면 솔직히 모른다고 말합니다. 여기서 Assistant의 이름은 서연입니다. 
        
            <question>            
            {question}
            </question>

            Assistant:"""

    else:  # English
        if conv_type == "normal": # for General Conversation
            prompt_template = """\n\nHuman: Using the following conversation, answer friendly for the newest question. If you don't know the answer, just say that you don't know, don't try to make up an answer. You will be acting as a thoughtful advisor.

            <history>
            {history}
            </history>
            
            <question>            
            {input}
            </question>

            Assistant:"""

        elif conv_type=='qa':  # for RAG
            prompt_template = """\n\nHuman: Here is pieces of context, contained in <context> tags. Provide a concise answer to the question at the end. If you don't know the answer, just say that you don't know, don't try to make up an answer. 
            
            <context>
            {context}
            </context>

            Go directly into the main points without the preamble. Do not include any additional information like newline characters "\n" or character counts in the result.
                        
            <question>
            {question}
            </question>

            Assistant:"""

        elif conv_type=="translation": 
            prompt_template = """\n\nHuman: Here is an article, contained in <article> tags. Translate the article to Korean. Put it in <result> tags.
            
            <article>
            {input}
            </article>
                        
            Assistant:"""
        
        elif conv_type == "sentiment":  # for sentiment, input
            prompt_template = """\n\nHuman: Here is <example> review and extracted topics and sentiments as <result>.

            <example>
            The room was small but clean and comfortable. The front desk was really busy and the check-in line was long, but the staff were professional and very pleasant with each person they helped. We will stay there again.
            </example>

            <result>
            Cleanliness: Positive, 
            Service: Positive
            </result>

            <review>
            {input}
            </review>
            
            Assistant:"""

        elif conv_type == "pii":  # removing PII(personally identifiable information) containing name, phone number, address
            prompt_template = """\n\nHuman: We want to de-identify some text by removing all personally identifiable information from this text so that it can be shared safely with external contractors.
            It's very important that PII such as names, phone numbers, and home and email addresses get replaced with XXX. Put it in <result> tags.

            Here is the text, inside <text></text> XML tags.

            <text>
            {input}
            </text>

            Assistant:"""

        elif conv_type == "extraction":  # for sentiment, input
            prompt_template = """\n\nHuman: Please precisely copy any email addresses from the following text and then write them, one per line.  Only write an email address if it's precisely spelled out in the input text.  If there are no email addresses in the text, write "N/A".  Do not say anything else.  Put it in <result> tags.

            {input}

            Assistant:"""

        elif conv_type == "grammar":  # Checking Grammatical Errors
            prompt_template = """\n\nHuman: Here is an article, contained in <article> tags:

            <article>
            {input}
            </article>

            Please identify any grammatical errors in the article. Also, add the fixed article at the end of answer.
            
            Assistant: """

        elif conv_type == "step-by-step":  # compelex question 
            prompt_template = """\n\nHuman: Using the following conversation, answer friendly for the newest question. If you don't know the answer, just say that you don't know, don't try to make up an answer. You will be acting as a thoughtful advisor.
            
            {input}

            Assistant: Can I think step by step?

            Human: Yes, please do.

            Assistant:"""
        
        elif conv_type == "like-child":  # Child Conversation (few shot)
            prompt_template = """\n\nHuman: Please complete the conversation by writing the next line, speaking as "A". You will be acting as a kindergarten teacher.

            Q: Is the tooth fairy real?
            A: Of course, sweetie. Wrap up your tooth and put it under your pillow tonight. There might be something waiting for you in the morning.
            Q: {input}

            Assistant:"""       

        elif conv_type == "funny": # for free conversation
            prompt_template = """\n\nHuman: 다음의 <history>는 Human과 Assistant의 친근한 이전 대화입니다. Assistant의 이름은 서서이고 10살 여자 어린이입니다. 상상력이 풍부하고 재미있는 대화를 잘합니다. 때론 바보같은 답변을 합니다.

            <history>
            {history}
            </history>

            <question>            
            {input}
            </question>
            
            Assistant:"""     

        else: # normal
            prompt_template = """\n\nHuman: Using the following conversation, answer friendly for the newest question. If you don't know the answer, just say that you don't know, don't try to make up an answer. You will be acting as a thoughtful advisor named Seoyeon.

            Human: {input}

            Assistant:"""
    
    return PromptTemplate.from_template(prompt_template)

# load documents from s3 
def load_document(file_type, s3_file_name):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, s3_prefix+'/'+s3_file_name)
    
    if file_type == 'pdf':
        Byte_contents = doc.get()['Body'].read()
        reader = PyPDF2.PdfReader(BytesIO(Byte_contents))
        
        texts = []
        for page in reader.pages:
            texts.append(page.extract_text())
        contents = '\n'.join(texts)
        
    elif file_type == 'pptx':
        Byte_contents = doc.get()['Body'].read()
            
        from pptx import Presentation
        prs = Presentation(BytesIO(Byte_contents))

        texts = []
        for i, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = text + shape.text
            texts.append(text)
        contents = '\n'.join(texts)
        
    elif file_type == 'txt' or file_type == 'md':        
        contents = doc.get()['Body'].read().decode('utf-8')

    elif file_type == 'docx':
        Byte_contents = doc.get()['Body'].read()
            
        import docx
        doc_contents =docx.Document(BytesIO(Byte_contents))

        texts = []
        for i, para in enumerate(doc_contents.paragraphs):
            if(para.text):
                texts.append(para.text)
                # print(f"{i}: {para.text}")        
        contents = '\n'.join(texts)
            
    # print('contents: ', contents)
    new_contents = str(contents).replace("\n"," ") 
    print('length: ', len(new_contents))

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=100,
        separators=["\n\n", "\n", ".", " ", ""],
        length_function = len,
    ) 

    texts = text_splitter.split_text(new_contents) 
                
    return texts

# load csv documents from s3
def load_csv_document(path, doc_prefix, s3_file_name):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, s3_prefix+'/'+s3_file_name)

    lines = doc.get()['Body'].read().decode('utf-8').split('\n')   # read csv per line
    print('lins: ', len(lines))
        
    columns = lines[0].split(',')  # get columns
    #columns = ["Category", "Information"]  
    #columns_to_metadata = ["type","Source"]
    print('columns: ', columns)
    
    docs = []
    n = 0
    for row in csv.DictReader(lines, delimiter=',',quotechar='"'):
        # print('row: ', row)
        #to_metadata = {col: row[col] for col in columns_to_metadata if col in row}
        values = {k: row[k] for k in columns if k in row}
        content = "\n".join(f"{k.strip()}: {v.strip()}" for k, v in values.items())
        doc = Document(
            page_content=content,
            metadata={
                'name': s3_file_name,
                'page': n+1,
                'uri': path+doc_prefix+parse.quote(s3_file_name)
            }
            #metadata=to_metadata
        )
        docs.append(doc)
        n = n+1
    print('docs[0]: ', docs[0])

    return docs

def get_summary(llm, texts):    
    # check korean
    pattern_hangul = re.compile('[\u3131-\u3163\uac00-\ud7a3]+') 
    word_kor = pattern_hangul.search(str(texts))
    print('word_kor: ', word_kor)
    
    if word_kor:
        prompt_template = """\n\nHuman: 다음 텍스트를 요약해서 500자 이내로 설명하세오.
        
        <text>
        {text}
        </text
        
        Assistant:"""        
    else:         
        prompt_template = """\n\nHuman: Write a concise summary of the following:

        {text}
        
        Assistant:"""
    
    PROMPT = PromptTemplate(template=prompt_template, input_variables=["text"])
    chain = load_summarize_chain(llm, chain_type="stuff", prompt=PROMPT)

    docs = [
        Document(
            page_content=t
        ) for t in texts[:5]
    ]
    summary = chain.run(docs)
    print('summary: ', summary)

    if summary == '':  # error notification
        summary = 'Fail to summarize the document. Try agan...'
        return summary
    else:
        # return summary[1:len(summary)-1]   
        return summary
    
def load_chat_history(userId, allowTime, conv_type, rag_type):
    dynamodb_client = boto3.client('dynamodb')

    response = dynamodb_client.query(
        TableName=callLogTableName,
        KeyConditionExpression='user_id = :userId AND request_time > :allowTime',
        ExpressionAttributeValues={
            ':userId': {'S': userId},
            ':allowTime': {'S': allowTime}
        }
    )
    # print('query result: ', response['Items'])

    for item in response['Items']:
        text = item['body']['S']
        msg = item['msg']['S']
        type = item['type']['S']

        if type == 'text':
            memory_chain.chat_memory.add_user_message(text)
            if len(msg) > MSG_LENGTH:
                memory_chain.chat_memory.add_ai_message(msg[:MSG_LENGTH])                          
            else:
                memory_chain.chat_memory.add_ai_message(msg) 
                
            if len(msg) > MSG_LENGTH:
                memory_chat.save_context({"input": text}, {"output": msg[:MSG_LENGTH]})
            else:
                memory_chat.save_context({"input": text}, {"output": msg})
                
def getAllowTime():
    d = datetime.datetime.now() - datetime.timedelta(days = 2)
    timeStr = str(d)[0:19]
    print('allow time: ',timeStr)

    return timeStr

def isTyping(connectionId, requestId):    
    msg_proceeding = {
        'request_id': requestId,
        'msg': 'Proceeding...',
        'status': 'istyping'
    }
    #print('result: ', json.dumps(result))
    sendMessage(connectionId, msg_proceeding)

def readStreamMsg(connectionId, requestId, stream):
    msg = ""
    if stream:
        for event in stream:
            #print('event: ', event)
            msg = msg + event

            result = {
                'request_id': requestId,
                'msg': msg,
                'status': 'proceeding'
            }
            #print('result: ', json.dumps(result))
            sendMessage(connectionId, result)
    # print('msg: ', msg)
    return msg

_ROLE_MAP = {"human": "\n\nHuman: ", "ai": "\n\nAssistant: "}
def extract_chat_history_from_memory():
    chat_history = []
    chats = memory_chain.load_memory_variables({})    
    # print('chats: ', chats)

    for dialogue_turn in chats['chat_history']:
        role_prefix = _ROLE_MAP.get(dialogue_turn.type, f"{dialogue_turn.type}: ")
        history = f"{role_prefix[2:]}{dialogue_turn.content}"
        if len(history)>MSG_LENGTH:
            chat_history.append(history[:MSG_LENGTH])
        else:
            chat_history.append(history)

    return chat_history

def get_revised_question(llm, connectionId, requestId, query):    
    # check korean
    pattern_hangul = re.compile('[\u3131-\u3163\uac00-\ud7a3]+')
    word_kor = pattern_hangul.search(str(query))
    print('word_kor: ', word_kor)

    if word_kor and word_kor != 'None':
        condense_template = """
        <history>
        {chat_history}
        </history>

        Human: <history>를 참조하여, 다음의 <question>의 뜻을 명확히 하는 새로운 질문을 한국어로 생성하세요. 새로운 질문은 원래 질문의 중요한 단어를 반드시 포함합니다.

        <question>            
        {question}
        </question>
            
        Assistant: 새로운 질문:"""
    else: 
        condense_template = """
        <history>
        {chat_history}
        </history>
        Answer only with the new question.

        Human: using <history>, rephrase the follow up <question> to be a standalone question. The standalone question must have main words of the original question.
         
        <quesion>
        {question}
        </question>

        Assistant: Standalone question:"""

    print('condense_template: ', condense_template)

    print('start prompt!')

    condense_prompt_claude = PromptTemplate.from_template(condense_template)
        
    condense_prompt_chain = LLMChain(llm=llm, prompt=condense_prompt_claude)

    chat_history = extract_chat_history_from_memory()
    try:         
        revised_question = condense_prompt_chain.run({"chat_history": chat_history, "question": query})
        print('revised_question: '+revised_question)
        
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                

        sendErrorMessage(connectionId, requestId, err_msg)        
        raise Exception ("Not able to request to LLM")
    
    if debugMessageMode=='true':
        # sendDebugMessage(connectionId, requestId, '[Debug]: '+revised_question)
        debug_msg_for_revised_question(llm, revised_question.replace("\n"," "), chat_history, connectionId, requestId)
    
    return revised_question.replace("\n"," ")

def priority_search(query, relevant_docs, bedrock_embedding, minSimilarity):
    excerpts = []
    for i, doc in enumerate(relevant_docs):
        # print('doc: ', doc)
        if 'translated_excerpt' in doc['metadata'] and doc['metadata']['translated_excerpt']:
            content = doc['metadata']['translated_excerpt']
        else:
            content = doc['metadata']['excerpt']
            
        print('content: ', content)
        
        excerpts.append(
            Document(
                page_content=content,
                metadata={
                    'name': doc['metadata']['title'],
                    'order':i,
                }
            )
        )  
    print('excerpts: ', excerpts)

    embeddings = bedrock_embedding
    vectorstore_confidence = FAISS.from_documents(
        excerpts,  # documents
        embeddings  # embeddings
    )            
    rel_documents = vectorstore_confidence.similarity_search_with_score(
        query=query,
        k=top_k
    )

    docs = []
    for i, document in enumerate(rel_documents):
        print(f'## Document(priority_search) {i+1}: {document}')

        order = document[0].metadata['order']
        name = document[0].metadata['name']
        assessed_score = document[1]
        print(f"{order} {name}: {assessed_score}")

        relevant_docs[order]['assessed_score'] = int(assessed_score)

        if assessed_score < minSimilarity:
            docs.append(relevant_docs[order])    
    # print('selected docs: ', docs)

    return docs

def get_reference(docs):
    reference = "\n\nFrom\n"
    for i, doc in enumerate(docs):
        if doc['metadata']['translated_excerpt']:
            excerpt = str(doc['metadata']['excerpt']+'  [번역]'+doc['metadata']['translated_excerpt']).replace('"',"") 
        else:
            excerpt = str(doc['metadata']['excerpt']).replace('"'," ")
            
        excerpt = excerpt.replace('\n','\\n')            
                
        if doc['rag_type'][:10] == 'opensearch':
            print(f'## Document(get_reference) {i+1}: {doc}')
                
            page = ""
            if "document_attributes" in doc['metadata']:
                if "_excerpt_page_number" in doc['metadata']['document_attributes']:
                    page = doc['metadata']['document_attributes']['_excerpt_page_number']
            uri = doc['metadata']['source']
            name = doc['metadata']['title']
            #print('opensearch page: ', page)

            if page:                
                reference = reference + f"{i+1}. {page}page in <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']})\n"
            else:
                reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
                    
        elif doc['rag_type'] == 'search': # google search
            print(f'## Document(get_reference) {i+1}: {doc}')
                
            uri = doc['metadata']['source']
            name = doc['metadata']['title']
            reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
                           
    return reference

def retrieve_docs_from_vectorstore(vectorstore_opensearch, query, top_k):
    print(f"query: {query}")

    relevant_docs = []
            
    # vector search (semantic) 
    relevant_documents = vectorstore_opensearch.similarity_search_with_score(
        query = query,
        k = top_k,
    )
    #print('(opensearch score) relevant_documents: ', relevant_documents)

    for i, document in enumerate(relevant_documents):
        #print('document.page_content:', document.page_content)
        #print('document.metadata:', document.metadata)
        print(f'## Document(opensearch-vector) {i+1}: {document}')

        name = document[0].metadata['name']
        # print('metadata: ', document[0].metadata)

        page = ""
        if "page" in document[0].metadata:
            page = document[0].metadata['page']
        uri = ""
        if "uri" in document[0].metadata:
            uri = document[0].metadata['uri']

        excerpt = document[0].page_content
        confidence = str(document[1])
        assessed_score = str(document[1])

        if page:
            print('page: ', page)
            doc_info = {
                "rag_type": 'opensearch-vector',
                "confidence": confidence,
                "metadata": {
                    "source": uri,
                    "title": name,
                    "excerpt": excerpt,
                    "translated_excerpt": "",
                    "document_attributes": {
                        "_excerpt_page_number": page
                    }
                },
                "assessed_score": assessed_score,
            }
        else:
            doc_info = {
                "rag_type": 'opensearch-vector',
                "confidence": confidence,
                "metadata": {
                    "source": uri,
                    "title": name,
                    "excerpt": excerpt,
                    "translated_excerpt": ""
                },
                "assessed_score": assessed_score,
            }
        relevant_docs.append(doc_info)
        
    return relevant_docs

def debug_msg_for_revised_question(llm, revised_question, chat_history, connectionId, requestId):
    global history_length, token_counter_history # debugMessageMode 

    history_context = ""
    token_counter_history = 0
    for history in chat_history:
        history_context = history_context + history

    if history_context:
        token_counter_history = llm.get_num_tokens(history_context)
        print('token_size of history: ', token_counter_history)

    history_length = len(history_context)

    sendDebugMessage(connectionId, requestId, f"새로운 질문: {revised_question}\n * 대화이력({str(history_length)}자, {token_counter_history} Tokens)을 활용하였습니다.")

def get_answer_using_RAG(llm, text, conv_type, connectionId, requestId, bedrock_embedding, rag_type):
    global time_for_revise, time_for_rag, time_for_inference, time_for_priority_search, number_of_relevant_docs  # for debug
    time_for_revise = time_for_rag = time_for_inference = time_for_priority_search = number_of_relevant_docs = 0

    global time_for_rag_inference
    time_for_rag_inference = 0
    
    vectorstore_opensearch = OpenSearchVectorSearch(
        index_name = "idx-*", # all
        is_aoss = False,
        ef_search = 1024, # 512(default)
        m=48,
        #engine="faiss",  # default: nmslib
        embedding_function = bedrock_embedding,
        opensearch_url=opensearch_url,
        http_auth=(opensearch_account, opensearch_passwd), # http_auth=awsauth,
    ) 
    
    reference = ""
    relevant_context = ""
    
    start_time_for_revise = time.time()

    revised_question = get_revised_question(llm, connectionId, requestId, text) # generate new prompt using chat history                    
    PROMPT = get_prompt_template(revised_question, conv_type, rag_type)
    # print('PROMPT: ', PROMPT)        

    end_time_for_revise = time.time()
    time_for_revise = end_time_for_revise - start_time_for_revise
    print('processing time for revised question: ', time_for_revise)

    relevant_docs = [] 
    rel_docs = retrieve_docs_from_vectorstore(vectorstore_opensearch=vectorstore_opensearch, query=revised_question, top_k=top_k)
    print(f'rel_docs ({rag_type}): '+json.dumps(rel_docs))
                
    if(len(rel_docs)>=1):
        for doc in rel_docs:
            relevant_docs.append(doc)

    if debugMessageMode=='true':
        for i, doc in enumerate(relevant_docs):
            print(f"#### relevant_docs ({i}): {json.dumps(doc)}")

    end_time_for_rag = time.time()
    time_for_rag = end_time_for_rag - end_time_for_revise
    print('processing time for RAG: ', time_for_rag)

    selected_relevant_docs = []
    if len(relevant_docs)>=1:
        print('start priority search')
        selected_relevant_docs = priority_search(revised_question, relevant_docs, bedrock_embedding, minDocSimilarity)
        print('selected_relevant_docs: ', json.dumps(selected_relevant_docs))

    if len(selected_relevant_docs)==0:
        print('No relevant document! So use google api')            
        api_key = google_api_key
        cse_id = google_cse_id 
            
        relevant_docs = []
        try: 
            service = build("customsearch", "v1", developerKey=api_key)
            result = service.cse().list(q=revised_question, cx=cse_id).execute()
            # print('google search result: ', result)

            if "items" in result:
                for item in result['items']:
                    api_type = "google api"
                    excerpt = item['snippet']
                    uri = item['link']
                    title = item['title']
                    confidence = ""
                    assessed_score = ""
                        
                    doc_info = {
                        "rag_type": 'search',
                        "api_type": api_type,
                        "confidence": confidence,
                        "metadata": {
                            "source": uri,
                            "title": title,
                            "excerpt": excerpt,
                            "translated_excerpt": "",
                        },
                        "assessed_score": assessed_score,
                    }
                    relevant_docs.append(doc_info)                
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)       

            sendErrorMessage(connectionId, requestId, "Not able to use Google API. Check the credentials")    
            #sendErrorMessage(connectionId, requestId, err_msg)    
            #raise Exception ("Not able to search using google api") 
            
        if len(relevant_docs)>=1:
            selected_relevant_docs = priority_search(revised_question, relevant_docs, bedrock_embedding, minDocSimilarity)
            print('selected_relevant_docs: ', json.dumps(selected_relevant_docs))
        # print('selected_relevant_docs (google): ', selected_relevant_docs)

        end_time_for_priority_search = time.time() 
        time_for_priority_search = end_time_for_priority_search - end_time_for_rag
        print('processing time for priority search: ', time_for_priority_search)
        number_of_relevant_docs = len(selected_relevant_docs)

        for document in selected_relevant_docs:
            if document['metadata']['translated_excerpt']:
                content = document['metadata']['translated_excerpt']
            else:
                content = document['metadata']['excerpt']
        
            relevant_context = relevant_context + content + "\n\n"
        print('relevant_context: ', relevant_context)

        try: 
            isTyping(connectionId, requestId)
            stream = llm(PROMPT.format(context=relevant_context, question=revised_question))
            msg = readStreamMsg(connectionId, requestId, stream)      
            # msg = msg.replace(" ","&nbsp;")        
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)       
            sendErrorMessage(connectionId, requestId, err_msg)    
            raise Exception ("Not able to request to LLM")    

        if len(selected_relevant_docs)>=1 and enableReference=='true':
            reference = get_reference(selected_relevant_docs)  

        end_time_for_inference = time.time()
        time_for_inference = end_time_for_inference - end_time_for_priority_search
        print('processing time for inference: ', time_for_inference)
    
    global relevant_length, token_counter_relevant_docs
    
    if debugMessageMode=='true':   # extract chat history for debug
        chat_history_all = extract_chat_history_from_memory()
        # print('chat_history_all: ', chat_history_all)

        history_length = []
        for history in chat_history_all:
            history_length.append(len(history))
        print('chat_history length: ', history_length)

        relevant_length = len(relevant_context)
        token_counter_relevant_docs = llm.get_num_tokens(relevant_context)

    memory_chain.chat_memory.add_user_message(text)  # append new diaglog
    memory_chain.chat_memory.add_ai_message(msg)

    return msg, reference

def get_answer_using_ConversationChain(text, conversation, conv_type, connectionId, requestId, rag_type):
    global time_for_inference
    
    start_time_for_inference = time.time()
    conversation.prompt = get_prompt_template(text, conv_type, rag_type)
    try: 
        isTyping(connectionId, requestId)    
        stream = conversation.predict(input=text)
        #print('stream: ', stream)                    
        msg = readStreamMsg(connectionId, requestId, stream)
        # msg = msg.replace(" ","&nbsp;")  
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")

    if debugMessageMode == 'true':   # extract chat history for debug
        chats = memory_chat.load_memory_variables({})
        chat_history_all = chats['history']
        # print('chat_history_all: ', chat_history_all)
        
        print('chat_history length: ', len(chat_history_all))
        
    end_time_for_inference = time.time()
    time_for_inference = end_time_for_inference - start_time_for_inference

    return msg

def get_answer_from_PROMPT(llm, text, conv_type, connectionId, requestId):
    PROMPT = get_prompt_template(text, conv_type, "")
    #print('PROMPT: ', PROMPT)

    try: 
        isTyping(connectionId, requestId)
        stream = llm(PROMPT.format(input=text))
        msg = readStreamMsg(connectionId, requestId, stream)
        # msg = msg.replace(" ","&nbsp;")  
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")
    
    return msg

def get_profile(function_type):
    claude3_sonnet = json.loads(os.environ.get('claude3_sonnet'))
    claude2 = json.loads(os.environ.get('claude2'))
    claude_instant = json.loads(os.environ.get('claude_instant'))

    if function_type == 'normal-claude3':
        profile_of_LLMs = claude3_sonnet
    elif function_type == 'normal-claude2':
        profile_of_LLMs = claude2
    elif function_type == 'normal-claude_instant':
        profile_of_LLMs = claude_instant
    else:
        profile_of_LLMs = json.loads(os.environ.get('profile_of_LLMs'))
            
    return profile_of_LLMs

def getResponse(connectionId, jsonBody):
    userId  = jsonBody['user_id']
    # print('userId: ', userId)
    requestId  = jsonBody['request_id']
    # print('requestId: ', requestId)
    requestTime  = jsonBody['request_time']
    # print('requestTime: ', requestTime)
    type  = jsonBody['type']
    # print('type: ', type)
    body = jsonBody['body']
    # print('body: ', body)
    conv_type = jsonBody['conv_type']  # conversation type
    print('Conversation Type: ', conv_type)
    function_type = jsonBody['function_type']  # conversation type
    print('Function Type: ', function_type)
    
    rag_type = ""
    if 'rag_type' in jsonBody:
        if jsonBody['rag_type']:
            rag_type = jsonBody['rag_type']  # RAG type
            print('rag_type: ', rag_type)

    global enableReference
    global map_chain, map_chat, memory_chat, memory_chain, debugMessageMode, selected_LLM, allowDualSearch
        
    profile_of_LLMs = get_profile(function_type)
    
    # Multi-region LLM
    if selected_LLM >= len(profile_of_LLMs)-1:
        selected_LLM = 0
        
    profile = profile_of_LLMs[selected_LLM]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'selected_LLM: {selected_LLM}, bedrock_region: {bedrock_region}, modelId: {modelId}')
    # print('profile: ', profile)
    
    llm = get_llm(profile_of_LLMs, selected_LLM)    
    bedrock_embedding = get_embedding(profile_of_LLMs, selected_LLM)

    # allocate memory
    if userId in map_chat or userId in map_chain:  
        print('memory exist. reuse it!')        
        memory_chain = map_chain[userId]
        memory_chat = map_chat[userId]
        conversation = ConversationChain(llm=llm, verbose=False, memory=memory_chat)
        
    else: 
        print('memory does not exist. create new one!')
        memory_chain = ConversationBufferWindowMemory(memory_key="chat_history", output_key='answer', return_messages=True, k=10)
        map_chain[userId] = memory_chain
        
        memory_chat = ConversationBufferWindowMemory(human_prefix='Human', ai_prefix='Assistant', k=MSG_HISTORY_LENGTH)
        map_chat[userId] = memory_chat

        allowTime = getAllowTime()
        load_chat_history(userId, allowTime, conv_type, rag_type)
        conversation = ConversationChain(llm=llm, verbose=False, memory=memory_chat)
 
    start = int(time.time())    

    msg = ""
    reference = ""
    token_counter_input = 0
    time_for_inference = 0
    history_length = 0
    isControlMsg = False
    
    if type == 'text' and body[:11] == 'list models':
        isControlMsg = True
        bedrock_client = boto3.client(
            service_name='bedrock',
            region_name=bedrock_region,
        )
        modelInfo = bedrock_client.list_foundation_models()    
        print('models: ', modelInfo)

        msg = f"The list of models: \n"
        lists = modelInfo['modelSummaries']
        
        for model in lists:
            msg += f"{model['modelId']}\n"
        
        msg += f"current model: {modelId}"
        print('model lists: ', msg)          

        sendResultMessage(connectionId, requestId, msg)  
    else:           
        text = body
        print('query: ', text)  
        querySize = len(text)
        textCount = len(text.split())
        print(f"query size: {querySize}, words: {textCount}")
        
        if type == 'text':
            if text == 'enableReference':
                enableReference = 'true'
                isControlMsg = True
                msg  = "Referece is enabled"
            elif text == 'disableReference':
                enableReference = 'false'
                isControlMsg = True
                msg  = "Reference is disabled"
            elif text == 'enableDebug':
                isControlMsg = True
                debugMessageMode = 'true'
                msg  = "Debug messages will be delivered to the client."
            elif text == 'disableDebug':
                isControlMsg = True
                debugMessageMode = 'false'
                msg  = "Debug messages will not be delivered to the client."

            elif text == 'clearMemory':
                memory_chain.clear()
                map_chain[userId] = memory_chain
                memory_chat.clear()                
                map_chat[userId] = memory_chat
                conversation = ConversationChain(llm=llm, verbose=False, memory=memory_chat)
                    
                print('initiate the chat memory!')
                msg  = "The chat memory was intialized in this session."
            else:       
                if conv_type == 'normal' or conv_type == 'funny':      # normal
                    msg = get_answer_using_ConversationChain(text, conversation, conv_type, connectionId, requestId, rag_type)
                
                elif conv_type == 'qa':   # RAG
                    print(f'rag_type: {rag_type}')
                    msg, reference = get_answer_using_RAG(llm, text, conv_type, connectionId, requestId, bedrock_embedding, rag_type)
                    
                elif conv_type == 'none':   # no prompt
                    try: 
                        msg = llm(HUMAN_PROMPT+text+AI_PROMPT)
                    except Exception:
                        err_msg = traceback.format_exc()
                        print('error message: ', err_msg)

                        sendErrorMessage(connectionId, requestId, err_msg)    
                        raise Exception ("Not able to request to LLM")
                else: 
                    msg = get_answer_from_PROMPT(llm, text, conv_type, connectionId, requestId)

                # token counter
                if debugMessageMode=='true':
                    token_counter_input = llm.get_num_tokens(text)
                    token_counter_output = llm.get_num_tokens(msg)
                    print(f"token_counter: question: {token_counter_input}, answer: {token_counter_output}")
                        
        elif type == 'document':
            isTyping(connectionId, requestId)
            
            object = body
            file_type = object[object.rfind('.')+1:len(object)]            
            print('file_type: ', file_type)
            
            if file_type == 'csv':
                docs = load_csv_document(path, doc_prefix, object)
                contexts = []
                for doc in docs:
                    contexts.append(doc.page_content)
                print('contexts: ', contexts)

                msg = get_summary(llm, contexts)
                        
            elif file_type == 'pdf' or file_type == 'txt' or file_type == 'md' or file_type == 'pptx' or file_type == 'docx':
                texts = load_document(file_type, object)

                docs = []
                for i in range(len(texts)):
                    docs.append(
                        Document(
                            page_content=texts[i],
                            metadata={
                                'name': object,
                                # 'page':i+1,
                                'uri': path+doc_prefix+parse.quote(object)
                            }
                        )
                    )
                print('docs[0]: ', docs[0])    
                print('docs size: ', len(docs))

                contexts = []
                for doc in docs:
                    contexts.append(doc.page_content)
                print('contexts: ', contexts)

                msg = get_summary(llm, contexts)
            else:
                msg = "uploaded file: "+object
                                                        
        sendResultMessage(connectionId, requestId, msg+reference)
        # print('msg+reference: ', msg+reference)

        elapsed_time = time.time() - start
        print("total run time(sec): ", elapsed_time)
                               
        item = {    # save dialog
            'user_id': {'S':userId},
            'request_id': {'S':requestId},
            'request_time': {'S':requestTime},
            'type': {'S':type},
            'body': {'S':body},
            'msg': {'S':msg+reference}
        }
        client = boto3.client('dynamodb')
        try:
            resp =  client.put_item(TableName=callLogTableName, Item=item)
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)
            raise Exception ("Not able to write into dynamodb")        
        #print('resp, ', resp)

    if debugMessageMode=='true' and isControlMsg==False: 
        statusMsg = f"\n[통계]\nRegion: {bedrock_region}\nModelId: {modelId}\n"
        if token_counter_input:
            statusMsg = statusMsg + f"Question: {str(len(text))}자 / {token_counter_input}토큰\nAnswer: {str(len(msg))}자 / {token_counter_output}토큰\n"
            
        if history_length:
            statusMsg = statusMsg + f"History: {str(history_length)}자 / {token_counter_history}토큰\n"
            
        statusMsg = statusMsg + f"Time(초): "            
        if time_for_inference != 0:
            statusMsg = statusMsg + f"{time_for_inference:.2f}(Inference), "
        statusMsg = statusMsg + f"{elapsed_time:.2f}(전체)"
            
        sendResultMessage(connectionId, requestId, msg+statusMsg)

    if selected_LLM >= len(profile_of_LLMs)-1:
        selected_LLM = 0
    else:
        selected_LLM = selected_LLM + 1

    return msg, reference

def lambda_handler(event, context):
    # print('event: ', event)
    
    msg = ""
    if event['requestContext']: 
        connectionId = event['requestContext']['connectionId']        
        routeKey = event['requestContext']['routeKey']
        
        if routeKey == '$connect':
            print('connected!')
        elif routeKey == '$disconnect':
            print('disconnected!')
        else:
            body = event.get("body", "")
            #print("data[0:8]: ", body[0:8])
            if body[0:8] == "__ping__":
                # print("keep alive!")                
                sendMessage(connectionId, "__pong__")
            else:
                print('connectionId: ', connectionId)
                print('routeKey: ', routeKey)
        
                jsonBody = json.loads(body)
                print('request body: ', json.dumps(jsonBody))

                requestId  = jsonBody['request_id']
                try:
                    msg, reference = getResponse(connectionId, jsonBody)

                    print('msg+reference: ', msg+reference)
                                        
                except Exception:
                    err_msg = traceback.format_exc()
                    print('err_msg: ', err_msg)

                    sendErrorMessage(connectionId, requestId, err_msg)    
                    raise Exception ("Not able to send a message")

    return {
        'statusCode': 200
    }
