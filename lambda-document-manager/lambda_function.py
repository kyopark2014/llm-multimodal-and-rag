import json
import boto3
import os
import traceback
import PyPDF2
import time
import docx

from io import BytesIO
from urllib import parse
from botocore.config import Config
from urllib.parse import unquote_plus
from langchain_community.embeddings import BedrockEmbeddings
from langchain_community.vectorstores.opensearch_vector_search import OpenSearchVectorSearch
from langchain_community.docstore.document import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from opensearchpy import OpenSearch
from pptx import Presentation
from multiprocessing import Process, Pipe
from langchain_core.prompts import ChatPromptTemplate
from langchain_community.chat_models import BedrockChat
from langchain.callbacks.streaming_stdout import StreamingStdOutCallbackHandler

s3 = boto3.client('s3')

s3_bucket = os.environ.get('s3_bucket') # bucket name
s3_prefix = os.environ.get('s3_prefix')
meta_prefix = "metadata/"
profile_of_LLMs = json.loads(os.environ.get('profile_of_LLMs'))
enableParallelSummay = os.environ.get('enableParallelSummay')

opensearch_account = os.environ.get('opensearch_account')
opensearch_passwd = os.environ.get('opensearch_passwd')
opensearch_url = os.environ.get('opensearch_url')
sqsUrl = os.environ.get('sqsUrl')
doc_prefix = s3_prefix+'/'

sqs = boto3.client('sqs')

roleArn = os.environ.get('roleArn') 
path = os.environ.get('path')
max_object_size = int(os.environ.get('max_object_size'))

supportedFormat = json.loads(os.environ.get('supportedFormat'))
print('supportedFormat: ', supportedFormat)

os_client = OpenSearch(
    hosts = [{
        'host': opensearch_url.replace("https://", ""), 
        'port': 443
    }],
    http_compress = True,
    http_auth=(opensearch_account, opensearch_passwd),
    use_ssl = True,
    verify_certs = True,
    ssl_assert_hostname = False,
    ssl_show_warn = False,
)

def delete_index_if_exist(index_name):    
    if os_client.indices.exists(index_name):
        print('delete opensearch document index: ', index_name)
        response = os_client.indices.delete(
            index=index_name
        )
        print('removed index: ', response)    
    else:
        print('no index: ', index_name)

# embedding for RAG
region_name = os.environ.get('bedrock_region')
    
boto3_bedrock = boto3.client(
    service_name='bedrock-runtime',
    region_name=region_name,
    config=Config(
        retries = {
            'max_attempts': 30
        }            
    )
)

bedrock_embeddings = BedrockEmbeddings(
    client=boto3_bedrock,
    region_name = region_name,
    model_id = 'amazon.titan-embed-text-v1' 
)   

def store_document_for_opensearch(bedrock_embeddings, docs, documentId):
    index_name = get_index_name(documentId)
    
    delete_index_if_exist(index_name)

    try:
        vectorstore = OpenSearchVectorSearch(
            index_name=index_name,  
            is_aoss = False,
            #engine="faiss",  # default: nmslib
            embedding_function = bedrock_embeddings,
            opensearch_url = opensearch_url,
            http_auth=(opensearch_account, opensearch_passwd),
        )
        response = vectorstore.add_documents(docs, bulk_size = 2000)
        print('response of adding documents: ', response)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                
        #raise Exception ("Not able to request to LLM")

    print('uploaded into opensearch')
    
def store_document_for_opensearch_with_nori(bedrock_embeddings, docs, documentId):
    index_name = get_index_name(documentId)
    
    delete_index_if_exist(index_name)
    
    index_body = {
        'settings': {
            'analysis': {
                'analyzer': {
                    'my_analyzer': {
                        'char_filter': ['html_strip'], 
                        'tokenizer': 'nori',
                        'filter': ['nori_number','lowercase','trim','my_nori_part_of_speech'],
                        'type': 'custom'
                    }
                },
                'tokenizer': {
                    'nori': {
                        'decompound_mode': 'mixed',
                        'discard_punctuation': 'true',
                        'type': 'nori_tokenizer'
                    }
                },
                "filter": {
                    "my_nori_part_of_speech": {
                        "type": "nori_part_of_speech",
                        "stoptags": [
                                "E", "IC", "J", "MAG", "MAJ",
                                "MM", "SP", "SSC", "SSO", "SC",
                                "SE", "XPN", "XSA", "XSN", "XSV",
                                "UNA", "NA", "VSV"
                        ]
                    }
                }
            },
            'index': {
                'knn': True,
                'knn.space_type': 'cosinesimil'  # Example space type
            }
        },
        'mappings': {
            'properties': {
                'metadata': {
                    'properties': {
                        'source' : {'type': 'keyword'},                    
                        'last_updated': {'type': 'date'},
                        'project': {'type': 'keyword'},
                        'seq_num': {'type': 'long'},
                        'title': {'type': 'text'},  # For full-text search
                        'url': {'type': 'text'},  # For full-text search
                    }
                },            
                'text': {
                    'analyzer': 'my_analyzer',
                    'search_analyzer': 'my_analyzer',
                    'type': 'text'
                },
                'vector_field': {
                    'type': 'knn_vector',
                    'dimension': 1536  # Replace with your vector dimension
                }
            }
        }
    }
    
    try: # create index
        response = os_client.indices.create(
            index_name,
            body=index_body
        )
        print('index was created with nori plugin:', response)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                
        #raise Exception ("Not able to create the index")

    try: # put the doucment
        vectorstore = OpenSearchVectorSearch(
            index_name=index_name,  
            is_aoss = False,
            #engine="faiss",  # default: nmslib
            embedding_function = bedrock_embeddings,
            opensearch_url = opensearch_url,
            http_auth=(opensearch_account, opensearch_passwd),
        )
        response = vectorstore.add_documents(docs, bulk_size = 2000)
        print('response of adding documents: ', response)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                
        #raise Exception ("Not able to request to LLM")

    print('uploaded into opensearch')    
 
def get_index_name(documentId):
    index_name = "idx-"+documentId
    # print('index_name: ', index_name)
                        
    print('index_name: ', index_name)
    print('length of index_name: ', len(index_name))
                            
    if len(index_name)>=100: # reduce index size
        index_name = 'idx-'+index_name[len(index_name)-100:]
        print('modified index_name: ', index_name)
    
    return index_name
 
# load documents from s3 for pdf and txt
def load_document(file_type, key):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, key)
    
    contents = ""
    if file_type == 'pdf':
        Byte_contents = doc.get()['Body'].read()
        
        try: 
            reader = PyPDF2.PdfReader(BytesIO(Byte_contents))
            
            texts = []
            for page in reader.pages:
                texts.append(page.extract_text())
            contents = '\n'.join(texts)
        except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to load the pdf file")
                
    elif file_type == 'pptx':
        Byte_contents = doc.get()['Body'].read()
            
        try:
            prs = Presentation(BytesIO(Byte_contents))

            texts = []
            for i, slide in enumerate(prs.slides):
                text = ""
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = text + shape.text
                texts.append(text)
            contents = '\n'.join(texts)
        except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to load texts from preseation file")
        
    elif file_type == 'txt' or file_type == 'md':       
        try:  
            contents = doc.get()['Body'].read().decode('utf-8')
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)        
            # raise Exception ("Not able to load the file")

    elif file_type == 'docx':
        try:
            Byte_contents = doc.get()['Body'].read()                    
            doc_contents =docx.Document(BytesIO(Byte_contents))

            texts = []
            for i, para in enumerate(doc_contents.paragraphs):
                if(para.text):
                    texts.append(para.text)
                    # print(f"{i}: {para.text}")        
            contents = '\n'.join(texts)            
            # print('contents: ', contents)
        except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to load docx")   
    
    texts = ""
    if len(contents)>0:
        new_contents = str(contents).replace("\n"," ") 
        print('length: ', len(new_contents))
        
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=100,
            separators=["\n\n", "\n", ".", " ", ""],
            length_function = len,
        ) 

        texts = text_splitter.split_text(new_contents) 
                        
    return texts

# load a code file from s3
def load_code(file_type, key):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, key)
    
    if file_type == 'py':        
        contents = doc.get()['Body'].read().decode('utf-8')
        separators = ["\ndef "]
        #print('contents: ', contents)
    elif file_type == 'js':
        contents = doc.get()['Body'].read().decode('utf-8')
        separators = ["\nfunction ", "\nexports.handler "]
    
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=50,
        chunk_overlap=0,
        #separators=["def ", "\n\n", "\n", ".", " ", ""],
        separators=separators,
        length_function = len,
    ) 

    texts = text_splitter.split_text(contents) 
    
    for i, text in enumerate(texts):
        print(f"Chunk #{i}: {text}")
                
    return texts

def isSupported(type):
    for format in supportedFormat:
        if type == format:
            return True
    
    return False
    
def check_supported_type(key, file_type, size):    
    if key.find('/html/') != -1 or key.find('/node_modules/') != -1 or key.find('/.git/') != -1: # do not include html/node_modules folders
        print('html: ', key.find('/html/'))
        return False
    
    if size > 5000 and size<max_object_size and isSupported(file_type):
        return True
    if size > 0 and (file_type == 'txt' or file_type == 'md' or file_type == 'py' or file_type == 'js'):
        return True
    elif size > 0 and (file_type == 'png' or file_type == 'jpeg' or file_type == 'jpg'):
        return True
    else:
        return False

HUMAN_PROMPT = "\n\nHuman:"
AI_PROMPT = "\n\nAssistant:"
def get_parameter(model_type, maxOutputTokens):
    if model_type=='titan': 
        return {
            "maxTokenCount":1024,
            "stopSequences":[],
            "temperature":0,
            "topP":0.9
        }
    elif model_type=='claude':
        return {
            "max_tokens_to_sample":maxOutputTokens, # 8k    
            "temperature":0.1,
            "top_k":250,
            "top_p":0.9,
            "stop_sequences": [HUMAN_PROMPT]            
        }
        
# Multi-LLM
def get_chat(profile_of_LLMs, selected_LLM):
    profile = profile_of_LLMs[selected_LLM]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'LLM: {selected_LLM}, bedrock_region: {bedrock_region}, modelId: {modelId}')
    maxOutputTokens = int(profile['maxOutputTokens'])
                          
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }            
        )
    )
    parameters = {
        "max_tokens":maxOutputTokens,     
        "temperature":0.1,
        "top_k":250,
        "top_p":0.9,
        "stop_sequences": [HUMAN_PROMPT]
    }
    # print('parameters: ', parameters)

    chat = BedrockChat(
        model_id=modelId,
        client=boto3_bedrock, 
        streaming=True,
        callbacks=[StreamingStdOutCallbackHandler()],
        model_kwargs=parameters,
    )        
    
    return chat

def summary_of_code(chat, code, mode):
    if mode == 'py': 
        system = (
            "다음의 <article> tag에는 python code가 있습니다. code의 전반적인 목적에 대해 설명하고, 각 함수의 기능과 역할을 자세하게 한국어 500자 이내로 설명하세요."
        )
    elif mode == 'js':
        system = (
            "다음의 <article> tag에는 node.js code가 있습니다. code의 전반적인 목적에 대해 설명하고, 각 함수의 기능과 역할을 자세하게 한국어 500자 이내로 설명하세요."
        )
    else:
        system = (
            "다음의 <article> tag에는 code가 있습니다. code의 전반적인 목적에 대해 설명하고, 각 함수의 기능과 역할을 자세하게 한국어 500자 이내로 설명하세요."
        )
    
    human = "<article>{code}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "code": code
            }
        )
        
        summary = result.content
        print('result of code summarization: ', summary)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return summary

def summarize_process_for_relevent_code(conn, chat, code, key, region_name):
    try: 
        if code.find('\ndef ') != -1:
            start = code.find('\ndef ')
            end = code.find(':')   
        elif code.find('\nfunction ') != -1:
            start = code.find('\nfunction ')
            end = code.find('(')   
        elif code.find('\nexports.') != -1:
            start = code.find('\nexports.')
            end = code.find(' =')         
        else:
            start = -1
            end = -1
              
        print('code: ', code)                             
        print(f'start: {start}, end: {end}')
                    
        doc = ""    
        if start != -1:      
            function_name = code[start+1:end]
            print('function_name: ', function_name)
            
            file_type = key[key.rfind('.')+1:len(key)].lower()
            print('file_type: ', file_type)
                            
            summary = summary_of_code(chat, code, file_type)
            print(f"summary ({region_name}, {file_type}): {summary}")
            
            # print('first line summary: ', summary[:len(function_name)])
            # print('function name: ', function_name)            
            if summary[:len(function_name)]==function_name:
                summary = summary[summary.find('\n')+1:len(summary)]

            doc = Document(
                page_content=summary,
                metadata={
                    'name': key,
                    # 'uri': path+doc_prefix+parse.quote(key),
                    'uri': path+key,
                    'code': code,
                    'function_name': function_name
                }
            )           
                        
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)       
        # raise Exception (f"Not able to summarize: {doc}")               
    
    conn.send(doc)    
    conn.close()

def summarize_relevant_codes_using_parallel_processing(codes, key):
    selected_LLM = 0
    relevant_codes = []    
    processes = []
    parent_connections = []
    for code in codes:
        parent_conn, child_conn = Pipe()
        parent_connections.append(parent_conn)
            
        chat = get_chat(profile_of_LLMs, selected_LLM)
        region_name = profile_of_LLMs[selected_LLM]['bedrock_region']

        process = Process(target=summarize_process_for_relevent_code, args=(child_conn, chat, code, key, region_name))
        processes.append(process)

        selected_LLM = selected_LLM + 1
        if selected_LLM == len(profile_of_LLMs):
            selected_LLM = 0

    for process in processes:
        process.start()
            
    for parent_conn in parent_connections:
        doc = parent_conn.recv()
        
        if doc:
            relevant_codes.append(doc)    

    for process in processes:
        process.join()
    
    return relevant_codes

def get_documentId(key, category):
    documentId = category + "-" + key
    documentId = documentId.replace(' ', '_') # remove spaces  
    documentId = documentId.replace(',', '_') # remove commas # not allowed: [ " * \\ < | , > / ? ]
    documentId = documentId.replace('/', '_') # remove slash
    documentId = documentId.lower() # change to lowercase
                
    return documentId
    
# load csv documents from s3
def lambda_handler(event, context):
    print('event: ', event)    
    
    documentIds = []
    for record in event['Records']:
        receiptHandle = record['receiptHandle']
        print("receiptHandle: ", receiptHandle)
        
        body = record['body']
        print("body: ", body)
        
        jsonbody = json.loads(body)        
        bucket = jsonbody['bucket']        
        # translate utf8
        key = unquote_plus(jsonbody['key']) # url decoding
        print('bucket: ', bucket)
        print('key: ', key)        
        eventName = jsonbody['type']
        
        start_time = time.time()      
        
        file_type = key[key.rfind('.')+1:len(key)].lower()
        print('file_type: ', file_type)
            
        if eventName == 'ObjectRemoved:Delete':
            if isSupported(file_type):
                objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
                print('objectName: ', objectName)
                
                # get metadata from s3
                metadata_key = meta_prefix+objectName+'.metadata.json'
                print('metadata_key: ', metadata_key)

                documentId = ""
                try: 
                    metadata_obj = s3.get_object(Bucket=bucket, Key=metadata_key)
                    metadata_body = metadata_obj['Body'].read().decode('utf-8')
                    metadata = json.loads(metadata_body)
                    print('metadata: ', metadata)
                    documentId = metadata['DocumentId']
                    print('documentId: ', documentId)
                    documentIds.append(documentId)
                except Exception:
                    err_msg = traceback.format_exc()
                    print('err_msg: ', err_msg)
                    # raise Exception ("Not able to get the object")
                    
                if documentId:
                    try: # delete metadata                        
                        print('delete metadata: ', metadata_key)                        
                        result = s3.delete_object(Bucket=bucket, Key=metadata_key)
                        # print('result of metadata deletion: ', result)
                        
                        # delete document index of opensearch
                        index_name = get_index_name(documentId)
                                                
                        delete_index_if_exist(index_name)                    
                    except Exception:
                        err_msg = traceback.format_exc()
                        print('err_msg: ', err_msg)
                        # raise Exception ("Not able to delete documents in Kendra")                    
            else: 
                print('This file format is not supported: ', file_type)                
                    
        elif eventName == "ObjectCreated:Put":            
            size = 0
            try:
                s3obj = s3.get_object(Bucket=bucket, Key=key)
                print(f"Got object: {s3obj}")        
                size = int(s3obj['ContentLength'])    
                
                #attributes = ['ETag', 'Checksum', 'ObjectParts', 'StorageClass', 'ObjectSize']
                #result = s3.get_object_attributes(Bucket=bucket, Key=key, ObjectAttributes=attributes)  
                #print('result: ', result)            
                #size = int(result['ObjectSize'])
                print('object size: ', size)
            except Exception:
                err_msg = traceback.format_exc()
                print('err_msg: ', err_msg)
                # raise Exception ("Not able to get object info") 
            
            if check_supported_type(key, file_type, size): 
                if file_type == 'py' or file_type == 'js':  # for code
                    category = file_type
                else:
                    category = "upload" # for document
                documentId = get_documentId(key, category)                                
                print('documentId: ', documentId)
                
                docs = []
                        
                if file_type == 'pdf' or file_type == 'txt' or file_type == 'md' or file_type == 'csv' or file_type == 'pptx' or file_type == 'docx':
                    print('upload to opensearch: ', key) 
                    texts = load_document(file_type, key)
                            
                    for i in range(len(texts)):
                        if texts[i]:
                            docs.append(
                                Document(
                                    page_content=texts[i],
                                    metadata={
                                        'name': key,
                                        # 'page':i+1,
                                        'uri': path+parse.quote(key)
                                    }
                                )
                            )
                                    
                elif file_type == 'py' or file_type == 'js':
                    codes = load_code(file_type, key)  # number of functions in the code
                                            
                    if enableParallelSummay=='true':
                        docs = summarize_relevant_codes_using_parallel_processing(codes, key)
                                
                    else:
                        for code in codes:
                            start = code.find('\ndef ')
                            end = code.find(':')                    
                            # print(f'start: {start}, end: {end}')
                                    
                            if start != -1:      
                                function_name = code[start+1:end]
                                # print('function_name: ', function_name)
                                                
                                chat = get_chat(profile_of_LLMs, 0)      
                                        
                                if file_type == 'py':
                                    mode = 'python'
                                elif file_type == 'js':
                                    mode = 'nodejs'
                                else:
                                    mode = file_type  
                                summary = summary_of_code(chat, code, mode)                        
                                            
                                if summary[:len(function_name)]==function_name:
                                    summary = summary[summary.find('\n')+1:len(summary)]
                                                                                        
                                docs.append(
                                    Document(
                                        page_content=summary,
                                        metadata={
                                            'name': key,
                                            # 'page':i+1,
                                            #'uri': path+doc_prefix+parse.quote(key),
                                            'uri': path+key,
                                            'code': code,
                                            'function_name': function_name
                                        }
                                    )
                                )                 
                                                                                                         
                print('docs size: ', len(docs))
                if len(docs)>0:
                    print('docs[0]: ', docs[0])
                                
                    store_document_for_opensearch(bedrock_embeddings, docs, documentId)

            else: # delete if the object is unsupported one for format or size
                try:
                    print('delete the unsupported file: ', key)                                
                    result = s3.delete_object(Bucket=bucket, Key=key)
                    print('result of deletion of the unsupported file: ', result)
                            
                except Exception:
                    err_msg = traceback.format_exc()
                    print('err_msg: ', err_msg)
                    # raise Exception ("Not able to delete unsupported file")
                    
        print('processing time: ', str(time.time() - start_time))
        
        # delete queue
        try:
            sqs.delete_message(QueueUrl=sqsUrl, ReceiptHandle=receiptHandle)
        except Exception as e:        
            print('Fail to delete the queue message: ', e)
            
    return {
        'statusCode': 200
    }
